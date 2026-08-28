from __future__ import annotations

import io
import zipfile

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md"}
ARCHIVE_EXTENSIONS = {".zip"}

_MAX_ZIP_MEMBERS = 20
_MAX_ZIP_UNCOMPRESSED_MB = 100


class ExtractionError(Exception):
    """Fichier illisible, corrompu, ou type non supporté."""


def _ext_of(filename: str) -> str:
    return "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def classify(filename: str) -> str:
    ext = _ext_of(filename)
    if ext in IMAGE_EXTENSIONS:
        return "image"
    if ext in DOCUMENT_EXTENSIONS or ext in ARCHIVE_EXTENSIONS:
        return "document"
    return "unknown"


def extract_text(filename: str, content: bytes) -> str:
    ext = _ext_of(filename)

    try:
        if ext == ".pdf":
            return _extract_pdf(content)
        if ext == ".docx":
            return _extract_docx(content)
        if ext == ".xlsx":
            return _extract_xlsx(content)
        if ext == ".pptx":
            return _extract_pptx(content)
        if ext in (".txt", ".md"):
            return content.decode("utf-8", errors="replace").strip()
        if ext == ".zip":
            return _extract_zip(content)
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(f"Fichier '{filename}' illisible ou corrompu : {exc}") from exc

    raise ExtractionError(f"Extraction non supportée pour '{ext}'")


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    return "\n".join(page.extract_text() or "" for page in reader.pages).strip()


def _extract_docx(content: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(content))
    return "\n".join(p.text for p in document.paragraphs).strip()


def _extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"# {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines).strip()


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
    return "\n".join(lines).strip()


def _extract_zip(content: bytes) -> str:
    try:
        zf = zipfile.ZipFile(io.BytesIO(content))
    except zipfile.BadZipFile as exc:
        raise ExtractionError(f"Archive zip corrompue : {exc}") from exc

    infos = [i for i in zf.infolist() if not i.is_dir()][:_MAX_ZIP_MEMBERS]
    total_uncompressed = sum(i.file_size for i in infos)
    if total_uncompressed > _MAX_ZIP_UNCOMPRESSED_MB * 1024 * 1024:
        raise ExtractionError("Archive zip trop volumineuse une fois décompressée.")

    parts: list[str] = []
    for info in infos:
        member_ext = _ext_of(info.filename)
        if member_ext not in DOCUMENT_EXTENSIONS:
            continue
        try:
            member_bytes = zf.read(info)
            member_text = extract_text(info.filename, member_bytes)
        except ExtractionError:
            continue
        if member_text:
            parts.append(f"## {info.filename}\n{member_text}")

    if not parts:
        raise ExtractionError("Aucun contenu exploitable dans l'archive zip.")
    return "\n\n".join(parts).strip()
